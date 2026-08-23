from __future__ import annotations

import asyncio
import hashlib
from importlib.metadata import version as package_version
import io
import json
import zipfile
from pathlib import Path
from threading import Event
from types import SimpleNamespace

import pytest
from fastapi.testclient import TestClient
from textual.widgets import RichLog
from typer.testing import CliRunner

from trisynapse_memory.api import create_app
from trisynapse_memory.cli import app
from trisynapse_memory.engine import MemoryEngine, SourceInput
from trisynapse_memory.engine.formation.sources import SourceError, _validate_public_url, prepare_source
from trisynapse_memory.engine.models import (
    IngestionRun,
    MemoryNamespace,
    QueryStep,
    SourceIngestionResult,
)
from trisynapse_memory.terminal import (
    COMPACT_LOGO,
    WIDE_LOGO,
    ModelSelectorScreen,
    MemoryTerminal,
    _path_suggestions,
    logo_for_width,
)


class Embedder:
    model_name = "v04-source-test"

    def encode(self, texts: list[str]) -> list[list[float]]:
        result = []
        for text in texts:
            vector = [0.0] * 16
            for token in text.lower().split():
                vector[int(hashlib.sha256(token.encode()).hexdigest(), 16) % 16] += 1
            norm = sum(value * value for value in vector) ** 0.5 or 1
            result.append([value / norm for value in vector])
        return result


def memory_at(path, **kwargs):
    return MemoryEngine.open(path, embedder=Embedder(), auto_process=False, **kwargs)


def test_code_directory_has_symbol_chunks_manifest_and_skips(tmp_path) -> None:
    repository = tmp_path / "repo"
    repository.mkdir()
    (repository / ".gitignore").write_text("ignored.py\n", encoding="utf-8")
    (repository / "ignored.py").write_text("def ignored(): pass\n", encoding="utf-8")
    (repository / ".env").write_text("TOKEN=secret\n", encoding="utf-8")
    (repository / "service.py").write_text(
        "import json\n\ndef answer(value: int) -> int:\n    return value + 1\n",
        encoding="utf-8",
    )

    engine = memory_at(tmp_path / "store")
    result = engine.ingest(SourceInput(kind="directory", path=str(repository), source_key="repo"))
    source = engine.get_source(result.source_id)
    deltas = [engine.get(delta_id) for delta_id in source.delta_ids]

    assert deltas[0].locator["kind"] == "manifest"
    symbol = next(item for item in deltas if item.locator["kind"] == "code_symbol")
    assert symbol.locator["symbol"] == "answer"
    assert symbol.locator["start_line"] == 3
    assert symbol.locator["end_line"] == 4
    assert "import json" in symbol.locator["metadata"]["imports"]
    assert symbol.payload["modality"] == "code"
    assert symbol.payload["retrieval_fields"]["symbol"] == "answer"
    assert symbol.payload["retrieval_fields"]["language"] == "py"
    assert "ignored.py" in source.metadata["skipped_paths"]
    assert ".env" in source.metadata["skipped_paths"]
    repeated = engine.ingest(SourceInput(kind="directory", path=str(repository), source_key="repo"))
    assert repeated.status == "skipped" and repeated.source_id == source.id


def test_mixed_ingestion_is_ordered_durable_and_failed_only_retry(tmp_path) -> None:
    missing = tmp_path / "later.md"
    engine = memory_at(tmp_path / "store")
    progress: list[str] = []
    run = engine.ingest_many([
        SourceInput(kind="text", text="A valid first source", source_key="first"),
        SourceInput(kind="file", path=str(missing), source_key="later"),
        SourceInput(kind="text", text="A valid third source", source_key="third"),
    ], on_progress=progress.append)

    assert run.status == "partial"
    assert [item.index for item in run.results] == [0, 1, 2]
    assert [item.status for item in run.results] == ["success", "failed", "success"]
    assert engine.get_ingestion_run(run.id).status == "partial"
    assert progress[0] == "Creating durable ingestion run"
    assert any(item.startswith("Preprocessed source") for item in progress)
    assert progress[-1] == "Finalizing ingestion run"

    missing.write_text("The failed item is available now.", encoding="utf-8")
    retried = engine.retry_ingestion(run.id)
    assert len(retried.inputs) == 1
    assert retried.results[0].status == "success"


def test_source_dedup_version_replacement_and_physical_removal(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    first = engine.ingest(SourceInput(kind="text", text="Version one", source_key="policy"))
    duplicate = engine.ingest(SourceInput(kind="text", text="Version one", source_key="policy"))
    second = engine.ingest(SourceInput(kind="text", text="Version two", source_key="policy"))

    assert duplicate.status == "skipped"
    assert duplicate.source_id == first.source_id
    first_record = engine.get_source(first.source_id)
    second_record = engine.get_source(second.source_id)
    assert first_record.status == "superseded"
    assert second_record.version == 2
    assert engine.get(first.delta_ids[0], include_retracted=True).text == "Version one"

    blob = engine.store.root / second_record.blob_path
    assert blob.is_file()
    removed = engine.remove_source(second_record.id, reason="test deletion")
    assert removed.removed_delta_ids == second_record.delta_ids
    assert engine.get(second_record.delta_ids[0], include_retracted=True).text == "[REMOVED]"
    assert not blob.exists()
    assert engine.validate_store().ok


def test_retained_original_is_in_backup_and_restore(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    result = engine.ingest(SourceInput(kind="text", text="retained original", source_key="backup-source"))
    source = engine.get_source(result.source_id)
    archive = engine.backup(tmp_path / "backup.zip")
    engine.close()

    restored = MemoryEngine.restore_backup(archive, tmp_path / "restored")
    restored_source = restored.get_source(source.id)
    assert (restored.store.root / restored_source.blob_path).read_text(encoding="utf-8") == "retained original"
    restored.close()


def test_blob_deduplication_is_content_only_and_reference_counted(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    first = engine.ingest(SourceInput(kind="text", text="same bytes", filename="one.txt", source_key="one"))
    second = engine.ingest(SourceInput(kind="text", text="same bytes", filename="two.md", source_key="two"))
    first_source = engine.get_source(first.source_id)
    second_source = engine.get_source(second.source_id)
    assert first_source.blob_path == second_source.blob_path
    blob = engine.store.root / first_source.blob_path
    engine.remove_source(first_source.id, reason="remove first reference")
    assert blob.exists()
    engine.remove_source(second_source.id, reason="remove last reference")
    assert not blob.exists()


def test_store_validation_detects_and_reingestion_repairs_corrupted_blob(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    first = engine.ingest(
        SourceInput(kind="text", text="authentic bytes", source_key="repairable")
    )
    source = engine.get_source(first.source_id)
    blob = engine.store.root / source.blob_path
    blob.write_bytes(b"corrupted bytes")

    invalid = engine.validate_store()
    assert invalid.ok is False
    assert source.id in invalid.corrupted_source_ids

    repeated = engine.ingest(
        SourceInput(kind="text", text="authentic bytes", source_key="repairable")
    )
    assert repeated.status == "skipped"
    assert blob.read_bytes() == b"authentic bytes"
    assert engine.validate_store().ok is True


def test_backup_refuses_a_store_with_a_corrupted_retained_source(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    result = engine.ingest(
        SourceInput(kind="text", text="original bytes", source_key="backup-check")
    )
    source = engine.get_source(result.source_id)
    (engine.store.root / source.blob_path).write_bytes(b"changed outside the engine")

    with pytest.raises(RuntimeError, match="failed validation before backup"):
        engine.backup(tmp_path / "invalid-backup.zip")

    assert not (tmp_path / "invalid-backup.zip").exists()


def test_archive_traversal_is_one_item_failure(tmp_path) -> None:
    payload = io.BytesIO()
    with zipfile.ZipFile(payload, "w") as archive:
        archive.writestr("../escape.md", "unsafe")
    engine = memory_at(tmp_path / "store")
    run = engine.ingest_many([
        SourceInput(kind="archive", content_base64=__import__("base64").b64encode(payload.getvalue()).decode(), filename="unsafe.zip"),
        SourceInput(kind="text", text="safe sibling", source_key="safe"),
    ])
    assert [item.status for item in run.results] == ["failed", "success"]
    assert "unsafe path" in run.results[0].error
    assert not (tmp_path / "escape.md").exists()


def test_image_uses_multimodal_provider_and_stores_extracted_text(tmp_path) -> None:
    class Vision:
        def __call__(self, system, user):
            return {}

        def complete_multimodal(self, system, user, image, media_type):
            assert "visible" in system.lower()
            assert media_type == "image/png"
            return {"description": "Screenshot shows api_key=super-secret-value", "visible_text": "Status OK"}

    engine = memory_at(tmp_path / "store", completion=Vision())
    encoded = __import__("base64").b64encode(b"not-a-real-image-but-provider-is-mocked").decode()
    result = engine.ingest(SourceInput(kind="image", content_base64=encoded, filename="status.png"))
    delta = engine.get(result.delta_ids[0])
    assert "super-secret-value" in delta.text
    assert "Status OK" in delta.text
    assert delta.privacy_scope == {}
    assert delta.payload["modality"] == "image"
    assert "Status OK" in delta.payload["retrieval_fields"]["visible_text"]


def test_source_rest_contract_and_breaking_remove_route(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    client = TestClient(create_app(tmp_path / "store", engine=engine))
    rejected = client.post("/api/v1/sources/ingest", json={"sources": [{"kind": "file", "path": "/tmp/server.txt"}]})
    assert rejected.status_code == 400

    accepted = client.post("/api/v1/sources/ingest", json={"sources": [{"kind": "text", "text": "REST source", "source_key": "rest"}]})
    assert accepted.status_code == 202
    run_id = accepted.json()["id"]
    run = client.get(f"/api/v1/ingestion-runs/{run_id}")
    assert run.status_code == 200
    source_id = run.json()["results"][0]["source_id"]
    assert client.get(f"/api/v1/sources/{source_id}").status_code == 200
    assert client.post("/api/v1/memory/purge", json={}).status_code == 404

    memory_id = run.json()["results"][0]["delta_ids"][0]
    removed = client.post("/api/v1/memory/remove", json={"delta_ids": [memory_id], "reason": "test", "confirm": True})
    assert removed.status_code == 200
    assert removed.json()["removed_delta_ids"] == [memory_id]
    source_record = engine.get_source(source_id)
    assert source_record.status == "active"
    assert (engine.store.root / source_record.blob_path).exists()


def test_office_and_notebook_locators(tmp_path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    docx_path = tmp_path / "guide.docx"
    document = Document()
    document.add_paragraph("First documented rule")
    document.save(docx_path)

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Release checklist"
    presentation.save(pptx_path)

    xlsx_path = tmp_path / "data.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Owners"
    sheet["A1"] = "Area"
    sheet["B1"] = "Owner"
    workbook.save(xlsx_path)

    notebook = tmp_path / "analysis.ipynb"
    notebook.write_text(json.dumps({"cells": [{"cell_type": "code", "source": ["answer = 42\n"]}]}), encoding="utf-8")

    engine = memory_at(tmp_path / "store")
    run = engine.ingest_many([SourceInput(kind="file", path=str(path)) for path in (docx_path, pptx_path, xlsx_path, notebook)])
    locators = [engine.get(result.delta_ids[0]).locator for result in run.results]
    assert locators[0]["kind"] == "paragraph"
    assert locators[1]["kind"] == "slide"
    assert locators[2]["sheet"] == "Owners" and locators[2]["row"] == 1
    assert locators[2]["cells"] == ["A1", "B1"]
    assert locators[3]["kind"] == "notebook_cell" and locators[3]["cell_index"] == 0
    modalities = [engine.get(result.delta_ids[0]).payload["modality"] for result in run.results]
    assert modalities == ["document", "document", "table", "code"]


def test_mocked_public_git_snapshot_records_commit_and_manifest(tmp_path, monkeypatch) -> None:
    from trisynapse_memory.engine.formation import sources as source_module

    def fake_run(command, **kwargs):
        if "clone" in command:
            target = Path(command[-1])
            target.mkdir(parents=True)
            (target / "service.py").write_text("def run():\n    return True\n", encoding="utf-8")
            return SimpleNamespace(stdout="")
        return SimpleNamespace(stdout="abc123def456\n")

    monkeypatch.setattr(source_module, "_validate_public_url", lambda *args, **kwargs: None)
    monkeypatch.setattr(source_module.subprocess, "run", fake_run)
    prepared = prepare_source(SourceInput(kind="git", url="https://github.com/example/public.git"))
    assert prepared.metadata["commit_sha"] == "abc123def456"
    assert prepared.metadata["remote_url"].startswith("https://")
    assert prepared.chunks[0].locator["kind"] == "manifest"
    assert any(chunk.locator.get("commit_sha") == "abc123def456" for chunk in prepared.chunks[1:])


def test_private_network_urls_and_images_without_vision_are_rejected(tmp_path, monkeypatch) -> None:
    from trisynapse_memory.engine.formation import sources as source_module

    monkeypatch.setattr(source_module.socket, "getaddrinfo", lambda *args, **kwargs: [(2, 1, 6, "", ("127.0.0.1", 80))])
    with pytest.raises(SourceError, match="private"):
        _validate_public_url("https://localhost/secret")

    encoded = __import__("base64").b64encode(b"image").decode()
    engine = memory_at(tmp_path / "store")
    run = engine.ingest_many([SourceInput(kind="image", content_base64=encoded, filename="image.png")])
    assert run.results[0].status == "failed"
    assert "vision support" in run.results[0].error


def test_legacy_chain_columns_and_purge_audit_are_removed_during_migration(tmp_path) -> None:
    engine = memory_at(tmp_path / "store")
    delta = engine.add("old sensitive content")
    connection = engine.store._connection
    connection.execute(
        "UPDATE deltas SET text='[PURGED]',payload_json=?,privacy_scope_json=? WHERE id=?",
        ('{"purged":true,"purge_id":"purge_old"}', '{"purged":true}', delta.id),
    )
    connection.execute("ALTER TABLE deltas ADD COLUMN prev_hash TEXT NOT NULL DEFAULT ''")
    connection.execute("ALTER TABLE deltas ADD COLUMN hash TEXT NOT NULL DEFAULT ''")
    connection.execute("ALTER TABLE snapshots ADD COLUMN evidence_hash TEXT NOT NULL DEFAULT ''")
    connection.execute("ALTER TABLE removal_audit ADD COLUMN old_root_hash TEXT NOT NULL DEFAULT ''")
    connection.execute("ALTER TABLE removal_audit ADD COLUMN new_root_hash TEXT NOT NULL DEFAULT ''")
    connection.execute("ALTER TABLE removal_audit RENAME TO purge_audit")
    connection.execute("ALTER TABLE purge_audit RENAME COLUMN remove_id TO purge_id")
    connection.execute(
        "INSERT INTO purge_audit VALUES(?,?,?,?,?,?,?)",
        (
            "purge_old",
            f'["{delta.id}"]',
            "legacy",
            "legacy deletion",
            "2025-01-01T00:00:00+00:00",
            "old",
            "new",
        ),
    )
    connection.commit()
    engine.close()

    reopened = memory_at(tmp_path / "store")
    migrated = reopened.get(delta.id, include_retracted=True)
    assert migrated.text == "[PURGED]"
    assert reopened.validate_store().ok
    tables = {row["name"] for row in reopened.store._connection.execute("SELECT name FROM sqlite_master WHERE type='table'")}
    assert "removal_audit" in tables and "purge_audit" not in tables
    delta_columns = {
        row["name"] for row in reopened.store._connection.execute("PRAGMA table_info(deltas)")
    }
    audit_columns = {
        row["name"] for row in reopened.store._connection.execute("PRAGMA table_info(removal_audit)")
    }
    snapshot_columns = {
        row["name"] for row in reopened.store._connection.execute("PRAGMA table_info(snapshots)")
    }
    assert "prev_hash" not in delta_columns and "hash" not in delta_columns
    assert "old_root_hash" not in audit_columns and "new_root_hash" not in audit_columns
    assert "evidence_hash" not in snapshot_columns


def test_cli_non_tty_help_version_and_json_contract() -> None:
    runner = CliRunner()
    no_args = runner.invoke(app, [])
    help_result = runner.invoke(app, ["--help"])
    version = runner.invoke(app, ["--version"])
    assert no_args.exit_code == 0
    assert "Store traces. Recall meaning." in no_args.output
    assert "●──" not in no_args.output
    assert "validate" in help_result.output
    assert "verify" not in help_result.output
    assert version.output.strip() == package_version("trisynapse-memory")
    assert logo_for_width(40) == COMPACT_LOGO
    assert logo_for_width(100) == WIDE_LOGO
    assert (len(WIDE_LOGO.splitlines()), max(map(len, WIDE_LOGO.splitlines()))) == (13, 52)
    assert (len(COMPACT_LOGO.splitlines()), max(map(len, COMPACT_LOGO.splitlines()))) == (8, 32)
    assert hashlib.sha256(WIDE_LOGO.encode()).hexdigest() == "f8f8b3b23b9bf266fdd33255a1881dc47568d843e49f3a96d4b3d341d061ce3a"
    assert hashlib.sha256(COMPACT_LOGO.encode()).hexdigest() == "b7e11da612e016b3f50809b4b1044c97bed3690521ffeaf4826ca0b94148e1a8"


def test_interactive_terminal_tabs_and_ctrl_c(tmp_path) -> None:
    async def exercise() -> None:
        terminal = MemoryTerminal(tmp_path / "terminal-store", MemoryNamespace())
        recent = terminal.engine.add("A memory ID completion target")
        retained = terminal.engine.ingest(
            SourceInput(kind="text", text="A retained terminal source", source_key="terminal")
        )

        def query_with_progress(*args, **kwargs):
            kwargs["on_step"](
                QueryStep(
                    id="q:test:classification",
                    phase="classification",
                    label="Classify query and prepare index",
                    sequence=1,
                    duration_ms=2,
                )
            )
            return SimpleNamespace(
                answer="The completion target is visible here.", citations=[]
            )

        terminal.engine.query = query_with_progress
        async with terminal.run_test(size=(100, 35)) as pilot:
            await pilot.pause()
            assert len(terminal.query("TabPane")) == 4
            assert COMPACT_LOGO in str(terminal.query_one("#brand").render())
            source_rows = "".join(
                line.text for line in terminal.query_one("#sources-log", RichLog).lines
            )
            inspector = terminal.query_one("TabbedContent")
            inspector.active = "trace-tab"
            await pilot.pause()
            trace_rows = "".join(
                line.text for line in terminal.query_one("#trace-log", RichLog).lines
            )
            assert retained.source_id in source_rows
            assert recent.id in trace_rows
            prompt = terminal.query_one("#prompt")
            inspector.active = "sources-tab"
            assert inspector.active == "sources-tab"
            prompt.value = "What is the completion target?"
            await pilot.press("enter")
            await pilot.pause()
            conversation = "\n".join(
                line.text for line in terminal.query_one("#activity", RichLog).lines
            )
            assert "What is the completion target?" in conversation
            assert "Classify query and prepare index" in conversation
            assert "The completion target is visible here." in conversation
            assert terminal._operation_label is None
            assert inspector.active == "sources-tab"
            prompt.value = "unfinished question"
            await pilot.press("ctrl+c")
            await pilot.pause()
            assert prompt.value == ""
            await pilot.press("/", "i", "n", "g")
            await pilot.pause()
            assert prompt._suggestion == "/ingest "
            await pilot.press("tab")
            await pilot.pause()
            assert prompt.value == "/ingest "
            assert "/ingest SOURCE" in str(terminal.query_one("#recommendations").render())
            prompt.value = ""
            suggested_memory_id = terminal._memory_ids()[0]
            await pilot.press("/", "h", "i", "s", "t", "o", "r", "y", " ", "d")
            await pilot.pause()
            assert prompt._suggestion == f"/history {suggested_memory_id}"
            await pilot.press("tab")
            assert prompt.value == f"/history {suggested_memory_id}"
            prompt.value = ""
            await pilot.press("/", "c", "o", "n", "f", "i", "g", "enter")
            await pilot.pause()
            assert terminal.query_one("TabbedContent").active == "config-tab"
            prompt.value = "/sources"
            await pilot.press("enter")
            await pilot.pause()
            assert terminal.query_one("TabbedContent").active == "sources-tab"
            source_log = terminal.query_one("#sources-log", RichLog)
            assert retained.source_id in "".join(line.text for line in source_log.lines)
            prompt.value = "/model embedding"
            await pilot.press("enter")
            await pilot.pause()
            assert isinstance(terminal.screen, ModelSelectorScreen)
            await pilot.click("#model-cancel")
            await pilot.pause()

    asyncio.run(exercise())


def test_ingestion_path_recommendations_follow_typed_characters(tmp_path, monkeypatch) -> None:
    (tmp_path / "source-code").mkdir()
    (tmp_path / "source-notes.md").write_text("notes", encoding="utf-8")
    monkeypatch.chdir(tmp_path)
    assert _path_suggestions("sou") == ["source-code/", "source-notes.md"]


def test_interactive_ingestion_keeps_the_input_responsive(tmp_path) -> None:
    async def exercise() -> None:
        terminal = MemoryTerminal(tmp_path / "responsive-store", MemoryNamespace())
        started = Event()
        release = Event()

        def slow_ingestion(sources, **kwargs):
            started.set()
            assert release.wait(timeout=5)
            return IngestionRun(
                id="ing_test",
                status="completed",
                inputs=sources,
                results=[
                    SourceIngestionResult(
                        index=0,
                        source_id="src_test",
                        kind="file",
                        status="success",
                    )
                ],
            )

        terminal.engine.ingest_many = slow_ingestion
        async with terminal.run_test(size=(100, 35)) as pilot:
            prompt = terminal.query_one("#prompt")
            prompt.value = "/ingest README.md"
            await pilot.press("enter")
            for _ in range(10):
                await pilot.pause()
                if started.is_set():
                    break
            assert started.is_set()
            assert "Reading and preprocessing" in str(
                terminal.query_one("#operation-status").render()
            )
            await pilot.press("x")
            assert prompt.value == "x"
            prompt.value = ""
            release.set()
            for _ in range(20):
                await pilot.pause()
                if terminal._operation_label is None:
                    break
            assert terminal._operation_label is None
            transcript = "\n".join(
                line.text for line in terminal.query_one("#activity", RichLog).lines
            )
            assert "Ingestion finished" in transcript

    asyncio.run(exercise())
