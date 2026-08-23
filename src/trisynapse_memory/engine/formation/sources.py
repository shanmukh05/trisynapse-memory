"""Unified, safety-bounded preprocessing for retained ingestion sources."""

from __future__ import annotations

import ast
import base64
import csv
import hashlib
import io
import ipaddress
import json
import mimetypes
import os
import re
import socket
import stat
import subprocess
import tarfile
import tempfile
import zipfile
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Any, Callable, Iterable
from urllib.parse import urlparse
from urllib.request import HTTPRedirectHandler, Request, build_opener

from trisynapse_memory.engine.models import SourceInput
from trisynapse_memory.prompts import load_prompt

MAX_FILE_BYTES = 25 * 1024 * 1024
MAX_RUN_BYTES = 250 * 1024 * 1024
MAX_ARCHIVE_FILES = 10_000
MAX_REDIRECTS = 5
FETCH_TIMEOUT = 30
CODE_SUFFIXES = {
    ".c", ".cc", ".cpp", ".cs", ".css", ".go", ".h", ".hpp", ".java",
    ".js", ".jsx", ".kt", ".lua", ".php", ".py", ".rb", ".rs", ".scala",
    ".sh", ".sql", ".swift", ".ts", ".tsx", ".vue", ".zig",
}
TEXT_SUFFIXES = {
    ".txt", ".md", ".mdx", ".rst", ".json", ".jsonl", ".yaml", ".yml",
    ".toml", ".csv",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp"}
ARCHIVE_SUFFIXES = {".zip", ".tar", ".tgz", ".tar.gz"}
OFFICE_SUFFIXES = {".docx", ".pptx", ".xlsx"}
HARD_IGNORES = {
    ".git", ".hg", ".svn", "node_modules", "vendor", "dist", "build", "target", ".next",
    ".cache", "__pycache__", ".venv", "venv", "coverage", ".pytest_cache", ".mypy_cache",
}
SECRET_NAMES = {
    ".env", ".env.local", ".env.production", "id_rsa", "id_ed25519", "credentials.json",
    "secrets.json", "service-account.json",
}
SECRET_SUFFIXES = {".pem", ".key", ".p12", ".pfx", ".keystore"}


@dataclass(frozen=True)
class LoadedDocument:
    text: str
    title: str
    media_type: str
    metadata: dict[str, Any] = field(default_factory=dict)


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._ignored = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        del attrs
        if tag in {"script", "style", "noscript"}:
            self._ignored += 1
        elif tag in {"p", "div", "section", "article", "h1", "h2", "h3", "h4", "li", "br", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self._ignored:
            self._ignored -= 1
        elif tag in {"p", "div", "section", "article", "li", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored:
            self.parts.append(data)


@dataclass(frozen=True)
class PreparedChunk:
    text: str
    locator: dict[str, Any]
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class PreparedSource:
    source: SourceInput
    kind: str
    title: str
    uri: str | None
    filename: str
    media_type: str
    original: bytes
    chunks: list[PreparedChunk]
    metadata: dict[str, Any] = field(default_factory=dict)
    skipped_paths: list[str] = field(default_factory=list)

    @property
    def content_hash(self) -> str:
        return hashlib.sha256(self.original).hexdigest()


class SourceError(ValueError):
    pass


def load_document(path: str | Path) -> LoadedDocument:
    """Load one local document into the legacy normalized-text contract."""

    file_path = Path(path).expanduser().resolve()
    if not file_path.is_file():
        raise FileNotFoundError(file_path)
    suffix = file_path.suffix.lower()
    media_type = mimetypes.guess_type(file_path.name)[0] or "application/octet-stream"
    metadata = {
        "path": str(file_path),
        "size_bytes": file_path.stat().st_size,
        "suffix": suffix,
    }
    if suffix == ".pdf":
        return _load_pdf_document(file_path, metadata)
    if suffix in {".html", ".htm"}:
        parser = _TextHTMLParser()
        parser.feed(file_path.read_text(encoding="utf-8", errors="replace"))
        text = "\n".join(
            line.strip()
            for line in "".join(parser.parts).splitlines()
            if line.strip()
        )
        return LoadedDocument(
            text=text,
            title=file_path.name,
            media_type="text/html",
            metadata=metadata,
        )
    if suffix in CODE_SUFFIXES:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        numbered = "\n".join(
            f"{index:06d}: {line}" for index, line in enumerate(raw.splitlines(), 1)
        )
        metadata["language"] = suffix.lstrip(".")
        return LoadedDocument(
            text=numbered,
            title=file_path.name,
            media_type=media_type,
            metadata=metadata,
        )
    if suffix in TEXT_SUFFIXES or media_type.startswith("text/"):
        return LoadedDocument(
            text=file_path.read_text(encoding="utf-8", errors="replace"),
            title=file_path.name,
            media_type=media_type,
            metadata=metadata,
        )
    raise ValueError(
        f"unsupported file type: {suffix or media_type}; "
        "supported: PDF, HTML, text, and source code"
    )


def _load_pdf_document(path: Path, metadata: dict[str, Any]) -> LoadedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError(
            "PDF loading requires: pip install 'trisynapse-memory[files]'"
        ) from exc
    reader = PdfReader(str(path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        pages.append(f"\n[Page {index}]\n{text.strip()}")
    metadata["page_count"] = len(reader.pages)
    title = str((reader.metadata or {}).get("/Title") or path.name)
    return LoadedDocument(
        text="\n".join(pages).strip(),
        title=title,
        media_type="application/pdf",
        metadata=metadata,
    )


def prepare_source(source: SourceInput, completion: Any | None = None) -> PreparedSource:
    if source.kind == "text":
        assert source.text is not None
        raw = source.text.encode("utf-8")
        return _text_source(source, raw, source.filename or "text.txt", source.title or "Text")
    if source.kind == "url":
        assert source.url is not None
        raw, final_url, media_type = _fetch(source.url)
        name = Path(urlparse(final_url).path).name or "page.html"
        return _bytes_source(source, raw, name, source.title or final_url, final_url, media_type, completion)
    if source.kind == "git":
        assert source.url is not None
        return _git_source(source, completion)

    if source.content_base64 is not None:
        try:
            raw = base64.b64decode(source.content_base64, validate=True)
        except ValueError as exc:
            raise SourceError("content_base64 is invalid") from exc
        if len(raw) > MAX_FILE_BYTES:
            raise SourceError("uploaded source exceeds the 25 MiB file limit")
        name = source.filename or "upload.bin"
        return _bytes_source(source, raw, name, source.title or name, None, None, completion)

    assert source.path is not None
    path = Path(source.path).expanduser().resolve()
    if path.is_symlink():
        raise SourceError("symlink sources are not accepted")
    if not path.exists():
        raise FileNotFoundError(path)
    if path.is_dir() or source.kind == "directory":
        return _directory_source(source, path, completion)
    raw = path.read_bytes()
    if len(raw) > MAX_FILE_BYTES:
        raise SourceError("file exceeds the 25 MiB source limit")
    return _bytes_source(source, raw, path.name, source.title or path.name, str(path), None, completion)


def store_blob(root: Path, content: bytes, *, filename: str) -> str:
    digest = hashlib.sha256(content).hexdigest()
    del filename  # Media type and original name live in SourceRecord, not the content address.
    directory = root / "sources" / "sha256" / digest[:2]
    directory.mkdir(parents=True, exist_ok=True, mode=0o700)
    try:
        directory.chmod(0o700)
    except OSError:
        pass
    target = directory / digest
    if (
        target.is_file()
        and not target.is_symlink()
        and hashlib.sha256(target.read_bytes()).hexdigest() == digest
    ):
        return str(target.relative_to(root))
    temporary: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="wb", dir=directory, prefix=f".{digest}.", delete=False
        ) as handle:
            temporary = Path(handle.name)
            handle.write(content)
            handle.flush()
            os.fsync(handle.fileno())
        temporary.replace(target)
        target.chmod(0o600)
    finally:
        if temporary is not None:
            temporary.unlink(missing_ok=True)
    return str(target.relative_to(root))


def _bytes_source(
    source: SourceInput,
    raw: bytes,
    filename: str,
    title: str,
    uri: str | None,
    media_type: str | None,
    completion: Any | None,
) -> PreparedSource:
    suffix = _suffix(filename)
    detected = media_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
    if suffix in IMAGE_SUFFIXES or source.kind == "image":
        return _image_source(source, raw, filename, title, uri, detected, completion)
    if _is_archive(filename) or source.kind == "archive":
        return _archive_source(source, raw, filename, title, uri, completion)
    chunks, metadata = _file_chunks(raw, filename)
    kind = "file"
    return PreparedSource(source, kind, title, uri, filename, detected, raw, chunks, {**source.metadata, **metadata})


def _text_source(source: SourceInput, raw: bytes, filename: str, title: str) -> PreparedSource:
    chunks = _chunk_text(raw.decode("utf-8", errors="replace"), filename)
    return PreparedSource(source, "text", title, None, filename, "text/plain", raw, chunks, source.metadata)


def _file_chunks(raw: bytes, filename: str) -> tuple[list[PreparedChunk], dict[str, Any]]:
    suffix = _suffix(filename)
    if suffix in CODE_SUFFIXES:
        text = raw.decode("utf-8", errors="replace")
        return _code_chunks(text, filename), {"language": suffix.lstrip("."), "source_type": "code"}
    if suffix == ".ipynb":
        return _notebook_chunks(raw, filename), {"source_type": "notebook"}
    if suffix == ".pdf":
        return _pdf_chunks(raw), {"source_type": "pdf"}
    if suffix == ".docx":
        return _docx_chunks(raw), {"source_type": "docx"}
    if suffix == ".pptx":
        return _pptx_chunks(raw), {"source_type": "pptx"}
    if suffix == ".xlsx":
        return _xlsx_chunks(raw), {"source_type": "xlsx"}
    if suffix in {".html", ".htm"}:
        parser = _TextHTMLParser()
        parser.feed(raw.decode("utf-8", errors="replace"))
        text = "\n".join(line.strip() for line in "".join(parser.parts).splitlines() if line.strip())
        return _chunk_text(text, filename), {"source_type": "html"}
    if suffix in {".json", ".jsonl", ".csv", ".yaml", ".yml"}:
        return _structured_chunks(raw, filename), {"source_type": "structured"}
    media_type = mimetypes.guess_type(filename)[0] or ""
    if suffix in TEXT_SUFFIXES or media_type.startswith("text/"):
        return _chunk_text(raw.decode("utf-8", errors="replace"), filename), {"source_type": "text"}
    raise SourceError(f"unsupported source type: {suffix or media_type or 'binary'}")


def _code_chunks(text: str, filename: str) -> list[PreparedChunk]:
    suffix = _suffix(filename)
    language = suffix.lstrip(".")
    lines = text.splitlines()
    spans: list[tuple[str, str, int, int]] = []
    if suffix == ".py":
        try:
            tree = ast.parse(text)
            for node in tree.body:
                if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                    spans.append((node.name, node.__class__.__name__.replace("Def", "").lower(), node.lineno, getattr(node, "end_lineno", node.lineno)))
        except SyntaxError:
            pass
    if not spans:
        spans = _tree_sitter_spans(text, language)
    if not spans:
        pattern = re.compile(r"^\s*(?:export\s+)?(?:async\s+)?(?:def|class|function|interface|type|struct|enum|fn|func)\s+([A-Za-z_$][\w$]*)")
        starts = [(match.group(1), index) for index, line in enumerate(lines, 1) if (match := pattern.match(line))]
        for pos, (name, start) in enumerate(starts):
            end = starts[pos + 1][1] - 1 if pos + 1 < len(starts) else len(lines)
            spans.append((name, "symbol", start, end))
    imports = [line.strip() for line in lines if re.match(r"^\s*(?:import|from|require\(|use\s|#include)", line)][:40]
    chunks: list[PreparedChunk] = []
    covered: set[int] = set()
    for name, symbol_kind, start, end in spans:
        end = min(end, len(lines))
        covered.update(range(start, end + 1))
        numbered = "\n".join(f"{line_no:06d}: {lines[line_no - 1]}" for line_no in range(start, end + 1))
        chunks.append(PreparedChunk(numbered, {
            "kind": "code_symbol", "path": filename, "language": language, "symbol": name,
            "symbol_kind": symbol_kind, "start_line": start, "end_line": end,
        }, {"imports": imports, "file_hash": hashlib.sha256(text.encode()).hexdigest()}))
    remaining = [index for index in range(1, len(lines) + 1) if index not in covered]
    for start in range(0, len(remaining), 120):
        group = remaining[start:start + 120]
        if not group:
            continue
        numbered = "\n".join(f"{line_no:06d}: {lines[line_no - 1]}" for line_no in group)
        chunks.append(PreparedChunk(numbered, {
            "kind": "code_lines", "path": filename, "language": language,
            "start_line": group[0], "end_line": group[-1],
        }, {"imports": imports, "file_hash": hashlib.sha256(text.encode()).hexdigest()}))
    return chunks or [PreparedChunk("", {"kind": "code_lines", "path": filename, "start_line": 1, "end_line": 0})]


def _tree_sitter_spans(text: str, language: str) -> list[tuple[str, str, int, int]]:
    aliases = {"js": "javascript", "jsx": "javascript", "ts": "typescript", "tsx": "tsx", "py": "python", "rb": "ruby", "sh": "bash", "cs": "c_sharp", "cpp": "cpp", "cc": "cpp", "h": "c"}
    try:
        from tree_sitter_language_pack import get_parser
        parser = get_parser(aliases.get(language, language))
        tree = parser.parse(text.encode("utf-8"))
    except Exception:
        return []
    interesting = {"function_definition", "function_declaration", "method_definition", "method_declaration", "class_definition", "class_declaration", "interface_declaration", "type_alias_declaration", "struct_item", "enum_item", "function_item", "impl_item"}
    result: list[tuple[str, str, int, int]] = []
    stack = [tree.root_node]
    while stack:
        node = stack.pop()
        if node.type in interesting:
            name_node = node.child_by_field_name("name")
            name = text.encode("utf-8")[name_node.start_byte:name_node.end_byte].decode("utf-8", errors="replace") if name_node else node.type
            result.append((name, node.type, node.start_point[0] + 1, node.end_point[0] + 1))
            continue
        stack.extend(reversed(node.children))
    return sorted(result, key=lambda item: item[2])


def _notebook_chunks(raw: bytes, filename: str) -> list[PreparedChunk]:
    payload = json.loads(raw.decode("utf-8"))
    chunks = []
    for index, cell in enumerate(payload.get("cells") or []):
        text = "".join(cell.get("source") or []).strip()
        if text:
            chunks.append(PreparedChunk(text, {"kind": "notebook_cell", "path": filename, "cell_index": index, "cell_type": cell.get("cell_type", "unknown")}))
    return chunks


def _pdf_chunks(raw: bytes) -> list[PreparedChunk]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise SourceError("PDF loading requires the sources extra") from exc
    reader = PdfReader(io.BytesIO(raw))
    chunks = []
    for index, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(PreparedChunk(text, {"kind": "page", "page": index}))
    return chunks


def _docx_chunks(raw: bytes) -> list[PreparedChunk]:
    try:
        from docx import Document
    except ImportError as exc:
        raise SourceError("DOCX loading requires the sources extra") from exc
    document = Document(io.BytesIO(raw))
    return [PreparedChunk(p.text.strip(), {"kind": "paragraph", "paragraph": index}) for index, p in enumerate(document.paragraphs, 1) if p.text.strip()]


def _pptx_chunks(raw: bytes) -> list[PreparedChunk]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SourceError("PPTX loading requires the sources extra") from exc
    presentation = Presentation(io.BytesIO(raw))
    chunks = []
    for slide_no, slide in enumerate(presentation.slides, 1):
        text = "\n".join(shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
        if text:
            chunks.append(PreparedChunk(text, {"kind": "slide", "slide": slide_no}))
    return chunks


def _xlsx_chunks(raw: bytes) -> list[PreparedChunk]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise SourceError("XLSX loading requires the sources extra") from exc
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    chunks = []
    for sheet in workbook.worksheets:
        for row_no, row in enumerate(sheet.iter_rows(), 1):
            values = ["" if cell.value is None else str(cell.value) for cell in row]
            if any(values):
                chunks.append(PreparedChunk(" | ".join(values), {
                    "kind": "sheet_row", "sheet": sheet.title, "row": row_no,
                    "cells": [cell.coordinate for cell in row if cell.value is not None],
                }))
    return chunks


def _structured_chunks(raw: bytes, filename: str) -> list[PreparedChunk]:
    suffix = _suffix(filename)
    text = raw.decode("utf-8", errors="replace")
    if suffix == ".jsonl":
        return [PreparedChunk(line, {"kind": "record", "record": index}) for index, line in enumerate(text.splitlines(), 1) if line.strip()]
    if suffix == ".csv":
        rows = csv.reader(io.StringIO(text))
        return [PreparedChunk(" | ".join(row), {"kind": "row", "row": index}) for index, row in enumerate(rows, 1)]
    if suffix == ".json":
        value = json.loads(text)
        values = value if isinstance(value, list) else [value]
        return [PreparedChunk(json.dumps(item, ensure_ascii=False, sort_keys=True), {"kind": "record", "record": index}) for index, item in enumerate(values, 1)]
    try:
        import yaml
    except ImportError:
        return _chunk_text(text, filename)
    values = list(yaml.safe_load_all(text))
    return [PreparedChunk(json.dumps(item, ensure_ascii=False, default=str, sort_keys=True), {"kind": "document", "document": index}) for index, item in enumerate(values, 1) if item is not None]


def _image_source(source: SourceInput, raw: bytes, filename: str, title: str, uri: str | None, media_type: str, completion: Any | None) -> PreparedSource:
    if completion is None or not hasattr(completion, "complete_multimodal"):
        raise SourceError("image ingestion requires a configured completion model with vision support")
    if getattr(completion, "vision_supported", None) is False:
        raise SourceError("the selected completion model is known not to support image input")
    prompt = load_prompt("image_extraction")
    payload = completion.complete_multimodal(
        prompt.text,
        "Extract a faithful memory-ready description of this image.",
        raw,
        media_type,
    )
    visible = str(payload.get("visible_text") or "").strip()
    description = str(payload.get("description") or "").strip()
    tables = payload.get("tables_or_charts") or []
    relationships = payload.get("relationships") or []
    text = "\n\n".join(part for part in [description, f"Visible text:\n{visible}" if visible else "", f"Tables/charts:\n{json.dumps(tables, ensure_ascii=False)}" if tables else "", f"Relationships:\n{json.dumps(relationships, ensure_ascii=False)}" if relationships else ""] if part)
    if not text:
        raise SourceError("vision model returned no usable image description")
    settings = getattr(completion, "settings", None)
    provenance = {
        "prompt": prompt.provenance(),
        "provider": getattr(settings, "provider", completion.__class__.__name__),
        "model": getattr(completion, "model", None),
    }
    return PreparedSource(
        source, "image", title, uri, filename, media_type, raw,
        [PreparedChunk(text, {"kind": "image", "filename": filename}, provenance)],
        {**source.metadata, "source_type": "image", "image_extraction": provenance},
    )


def _directory_source(source: SourceInput, root: Path, completion: Any | None, *, git_metadata: dict[str, Any] | None = None) -> PreparedSource:
    accepted, skipped = _collect_directory(root)
    if not accepted:
        raise SourceError("directory contains no supported source files")
    chunks: list[PreparedChunk] = []
    manifest = []
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as bundle:
        for path in accepted:
            relative = path.relative_to(root).as_posix()
            raw = path.read_bytes()
            entry = zipfile.ZipInfo(relative, date_time=(1980, 1, 1, 0, 0, 0))
            entry.compress_type = zipfile.ZIP_DEFLATED
            entry.external_attr = 0o600 << 16
            bundle.writestr(entry, raw)
            manifest.append(relative)
            if _suffix(relative) in IMAGE_SUFFIXES:
                image = _image_source(source, raw, relative, relative, None, mimetypes.guess_type(relative)[0] or "application/octet-stream", completion)
                file_chunks, metadata = image.chunks, image.metadata
            else:
                file_chunks, metadata = _file_chunks(raw, relative)
            for chunk in file_chunks:
                provenance = {key: value for key, value in (git_metadata or {}).items() if key in {"commit_sha", "remote_url"}}
                chunks.append(PreparedChunk(chunk.text, {**chunk.locator, "path": relative, **provenance}, {**metadata, **chunk.metadata}))
    manifest_text = "Repository/source manifest:\n" + "\n".join(manifest)
    chunks.insert(0, PreparedChunk(manifest_text, {"kind": "manifest", "path": "."}, {"file_count": len(manifest)}))
    metadata = {**source.metadata, **(git_metadata or {}), "file_count": len(manifest), "source_type": "repository"}
    title = source.title or root.name
    return PreparedSource(source, "directory" if not git_metadata else "git", title, str(root), f"{root.name}.zip", "application/zip", buffer.getvalue(), chunks, metadata, skipped)


def _git_source(source: SourceInput, completion: Any | None) -> PreparedSource:
    assert source.url is not None
    _validate_public_url(source.url, https_only=True)
    if not source.url.endswith(".git") and "github.com/" not in source.url and "gitlab.com/" not in source.url and "bitbucket.org/" not in source.url:
        raise SourceError("git source must be a public HTTPS repository URL")
    with tempfile.TemporaryDirectory(prefix="trisynapse-git-") as temp:
        target = Path(temp) / "repo"
        command = ["git", "clone", "--depth", "1", "--no-tags"]
        if source.ref:
            command.extend(["--branch", source.ref])
        command.extend([source.url, str(target)])
        try:
            subprocess.run(command, check=True, capture_output=True, text=True, timeout=120)
            commit = subprocess.run(["git", "-C", str(target), "rev-parse", "HEAD"], check=True, capture_output=True, text=True, timeout=10).stdout.strip()
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError) as exc:
            raise SourceError(f"public Git ingestion failed: {exc}") from exc
        titled_source = source if source.title else source.model_copy(
            update={"title": Path(urlparse(source.url).path).stem or "Git repository"}
        )
        prepared = _directory_source(titled_source, target, completion, git_metadata={"remote_url": source.url, "commit_sha": commit, "ref": source.ref})
        return PreparedSource(prepared.source, "git", prepared.title, source.url, prepared.filename, prepared.media_type, prepared.original, prepared.chunks, prepared.metadata, prepared.skipped_paths)


def _archive_source(source: SourceInput, raw: bytes, filename: str, title: str, uri: str | None, completion: Any | None) -> PreparedSource:
    with tempfile.TemporaryDirectory(prefix="trisynapse-archive-") as temp:
        root = Path(temp) / "content"
        root.mkdir()
        _extract_archive(raw, filename, root)
        prepared = _directory_source(source, root, completion)
        return PreparedSource(prepared.source, "archive", title, uri, filename, "application/zip", raw, prepared.chunks, {**prepared.metadata, "source_type": "archive"}, prepared.skipped_paths)


def _collect_directory(root: Path) -> tuple[list[Path], list[str]]:
    patterns = _ignore_patterns(root)
    accepted: list[Path] = []
    skipped: list[str] = []
    total = 0
    for path in sorted(root.rglob("*")):
        relative = path.relative_to(root).as_posix()
        parts = set(PurePosixPath(relative).parts)
        if path.is_symlink() or parts & HARD_IGNORES or path.name.lower() in SECRET_NAMES or path.suffix.lower() in SECRET_SUFFIXES or _matches_ignore(relative, patterns):
            skipped.append(relative)
            continue
        if not path.is_file():
            continue
        suffix = _suffix(path.name)
        media = mimetypes.guess_type(path.name)[0] or ""
        supported = suffix in CODE_SUFFIXES | TEXT_SUFFIXES | OFFICE_SUFFIXES | IMAGE_SUFFIXES | {".pdf", ".html", ".htm", ".ipynb"} or media.startswith("text/")
        if not supported or path.stat().st_size > MAX_FILE_BYTES:
            skipped.append(relative)
            continue
        total += path.stat().st_size
        if len(accepted) >= MAX_ARCHIVE_FILES or total > MAX_RUN_BYTES:
            raise SourceError("directory exceeds source file-count or expanded-size limit")
        accepted.append(path)
    return accepted, skipped


def _ignore_patterns(root: Path) -> list[str]:
    result = []
    for name in (".gitignore", ".trisynapseignore"):
        path = root / name
        if path.is_file():
            result.extend(line.strip() for line in path.read_text(encoding="utf-8", errors="replace").splitlines() if line.strip() and not line.lstrip().startswith("#"))
    return result


def _matches_ignore(relative: str, patterns: Iterable[str]) -> bool:
    from fnmatch import fnmatch

    ignored = False
    for pattern in patterns:
        negated = pattern.startswith("!")
        candidate = pattern[1:] if negated else pattern
        candidate = candidate.lstrip("/").rstrip("/")
        if fnmatch(relative, candidate) or fnmatch(relative, f"{candidate}/*"):
            ignored = not negated
    return ignored


def _extract_archive(raw: bytes, filename: str, root: Path) -> None:
    if filename.lower().endswith(".zip"):
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            members = []
            for item in archive.infolist():
                mode = item.external_attr >> 16
                if stat.S_ISLNK(mode):
                    raise SourceError("archive symlinks are not accepted")
                members.append((item.filename, item.file_size, item.is_dir(), lambda item=item: archive.open(item)))
            _write_archive_members(members, root)
        return
    mode = "r:gz" if filename.lower().endswith((".tgz", ".tar.gz")) else "r:"
    with tarfile.open(fileobj=io.BytesIO(raw), mode=mode) as archive:
        members = []
        for item in archive.getmembers():
            if item.issym() or item.islnk():
                raise SourceError("archive symlinks are not accepted")
            if item.isfile():
                members.append((item.name, item.size, False, lambda item=item: archive.extractfile(item)))
        _write_archive_members(members, root)


def _write_archive_members(members: Iterable[tuple[str, int, bool, Callable[[], Any]]], root: Path) -> None:
    total = 0
    count = 0
    for name, size, is_dir, opener in members:
        if is_dir:
            continue
        count += 1
        total += size
        if count > MAX_ARCHIVE_FILES or total > MAX_RUN_BYTES:
            raise SourceError("archive exceeds file-count or expanded-size limit")
        relative = PurePosixPath(name)
        if relative.is_absolute() or ".." in relative.parts:
            raise SourceError("archive contains an unsafe path")
        target = (root / Path(*relative.parts)).resolve()
        if root.resolve() not in target.parents:
            raise SourceError("archive contains an unsafe path")
        target.parent.mkdir(parents=True, exist_ok=True)
        handle = opener()
        if handle is None:
            continue
        with handle:
            target.write_bytes(handle.read())


class _SafeRedirect(HTTPRedirectHandler):
    def __init__(self) -> None:
        super().__init__()
        self.count = 0

    def redirect_request(self, req: Any, fp: Any, code: int, msg: str, headers: Any, newurl: str) -> Any:
        self.count += 1
        if self.count > MAX_REDIRECTS:
            raise SourceError("URL exceeded the redirect limit")
        _validate_public_url(newurl)
        return super().redirect_request(req, fp, code, msg, headers, newurl)


def _fetch(url: str) -> tuple[bytes, str, str]:
    _validate_public_url(url)
    opener = build_opener(_SafeRedirect())
    request = Request(url, headers={"User-Agent": "trisynapse-memory/0.4"})
    with opener.open(request, timeout=FETCH_TIMEOUT) as response:
        final = response.geturl()
        _validate_public_url(final)
        length = int(response.headers.get("Content-Length") or 0)
        if length > MAX_FILE_BYTES:
            raise SourceError("URL response exceeds the 25 MiB limit")
        raw = response.read(MAX_FILE_BYTES + 1)
        if len(raw) > MAX_FILE_BYTES:
            raise SourceError("URL response exceeds the 25 MiB limit")
        media_type = response.headers.get_content_type()
    return raw, final, media_type


def _validate_public_url(url: str, *, https_only: bool = False) -> None:
    parsed = urlparse(url)
    allowed = {"https"} if https_only else {"http", "https"}
    if parsed.scheme not in allowed or not parsed.hostname or parsed.username or parsed.password:
        raise SourceError("source URL must be a public HTTP(S) URL without embedded credentials")
    try:
        addresses = {item[4][0] for item in socket.getaddrinfo(parsed.hostname, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)}
    except socket.gaierror as exc:
        raise SourceError(f"source host could not be resolved: {parsed.hostname}") from exc
    for address in addresses:
        ip = ipaddress.ip_address(address)
        if not ip.is_global:
            raise SourceError("private, loopback, link-local, and reserved source addresses are blocked")


def _chunk_text(text: str, filename: str, *, chars: int = 3500) -> list[PreparedChunk]:
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n", text) if item.strip()]
    chunks: list[PreparedChunk] = []
    current = ""
    start = 1
    line = 1
    for paragraph in paragraphs:
        candidate = f"{current}\n\n{paragraph}".strip()
        if current and len(candidate) > chars:
            chunks.append(PreparedChunk(current, {"kind": "chunk", "path": filename, "start_line": start, "end_line": line - 1}))
            current = paragraph
            start = line
        else:
            current = candidate
        line += paragraph.count("\n") + 2
    if current:
        chunks.append(PreparedChunk(current, {"kind": "chunk", "path": filename, "start_line": start, "end_line": max(start, line - 1)}))
    return chunks


def _suffix(filename: str) -> str:
    lower = filename.lower()
    return ".tar.gz" if lower.endswith(".tar.gz") else Path(lower).suffix


def _is_archive(filename: str) -> bool:
    return _suffix(filename) in ARCHIVE_SUFFIXES
